# anki_suite.py — Tudo em um
# Reúne: kindle_to_md, txt2anki_api, md2anki_api, txt2anki_url e utilidades.
# Requisitos (dependendo das funções usadas):
#   pip install openai requests python-dotenv beautifulsoup4 python-pptx
#   (opcional) pip install spacy && python -m spacy download en_core_web_sm
#   Anki Desktop + AnkiConnect (porta 8765) aberto para envio de notas
#
# CLI:
#   python anki_suite.py kindle2md caminho/arquivo.html
#   python anki_suite.py txt2anki caminho/palavras.txt
#   python anki_suite.py md2anki caminho/arquivo.md [filtro.txt]
#   python anki_suite.py url2anki https://exemplo.com/estudo
#   python anki_suite.py write-ui caminho/saida/index.html
#
# Observação: Algumas funções usam recursos externos (OpenAI API, LibreOffice para PDF),
# portanto é normal que certas rotinas dependam do ambiente do usuário.

import os, re, sys, json, textwrap, random, hashlib, base64, tempfile, subprocess
from typing import List, Tuple, Optional, Dict, Any

# --- Imports de terceiros usados pelas rotinas ---
try:
    import requests
except Exception as e:
    print("⚠ A biblioteca 'requests' é necessária:", e)

try:
    from openai import OpenAI
except Exception as e:
    OpenAI = None
    print("⚠ A biblioteca 'openai' é necessária para certas rotinas:", e)

try:
    from bs4 import BeautifulSoup
except Exception as e:
    BeautifulSoup = None
    print("⚠ A biblioteca 'beautifulsoup4' é necessária para Kindle/URL:", e)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception as e:
    Presentation = None
    Inches = Pt = None
    print("⚠ A biblioteca 'python-pptx' é necessária para gerar relatório PPTX:", e)

# spaCy é opcional; se não estiver presente, fazemos fallback
try:
    import spacy
    _NLP = None
    def get_nlp():
        global _NLP
        if _NLP is None:
            _NLP = spacy.load("en_core_web_sm")
        return _NLP
except Exception:
    spacy = None
    def get_nlp():
        return None

# dotenv é opcional (carrega OPENAI_API_KEY do .env se existir)
try:
    import dotenv
    dotenv.load_dotenv()
except Exception:
    pass

# ============ Configuração comum (Anki/OpenAI) ============
CHAT_MODEL = os.getenv("CHAT_MODEL", "gpt-4o-mini")
TTS_MODEL  = os.getenv("TTS_MODEL", "tts-1-hd")
VOICES     = ["alloy","ash","coral","echo","fable","onyx","nova","shimmer","sage"]

DECK_NAME  = os.getenv("ANKI_DECK", "Teste de importação")
MODEL_NAME = os.getenv("ANKI_MODEL", "Supermemo-03572")
FIELDS     = ("Frente", "Verso")

ANKI_URL   = os.getenv("ANKI_URL", "http://127.0.0.1:8765")
MAX_TRIES  = 4

# Cliente OpenAI (se disponível)
_client = None
def get_client():
    global _client
    if _client is None:
        if OpenAI is None:
            raise RuntimeError("Pacote openai não encontrado.")
        _client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
    return _client

# ============ Utilitários Anki ============
def anki(action: str, **params):
    """Chama a API AnkiConnect; lança RuntimeError em caso de erro."""
    resp = requests.post(ANKI_URL, json={"action": action, "version": 6, "params": params}, timeout=60)
    data = resp.json()
    if data.get("error"):
        raise RuntimeError(str(data["error"]))
    return data.get("result")

def deck_setup():
    names = anki("deckNames") or []
    if DECK_NAME not in names:
        anki("createDeck", deck=DECK_NAME)
    models = anki("modelNames") or []
    if MODEL_NAME not in models:
        raise RuntimeError(f"Modelo “{MODEL_NAME}” não existe no Anki.")

def tts_mp3(text: str) -> str:
    """Gera um MP3 com voz aleatória e retorna o nome armazenado no Anki."""
    v = random.choice(VOICES)
    fn = f"tts_{v}_{hashlib.sha1(text.encode()).hexdigest()[:10]}.mp3"
    data = get_client().audio.speech.create(model=TTS_MODEL, voice=v, input=text).content
    anki("storeMediaFile", filename=fn, data=base64.b64encode(data).decode())
    return fn

def highlight_html(txt: str) -> str:
    return re.sub(r"«(.*?)»", r"<span style='color:#00aaff;font-weight:bold;'>\1</span>", txt)

PAIR = re.compile(r"S_([A-Z]+):\s*(.*?)\s*PT_\1:\s*(.*?)\s*(?=S_[A-Z]+:|$)", re.S)

# ============ kindle_to_md ============
# Baseado no seu kindle_to_md.py
_KINDLE_META = re.compile(r'\b(Destaque|Highlight)\b.*$', re.I)

def kindle_extract_html_to_md(html_text: str) -> str:
    """Extrai destaques de HTML do Kindle em Markdown (blocos de citação)."""
    if BeautifulSoup is None:
        raise RuntimeError("beautifulsoup4 não está instalado.")
    soup = BeautifulSoup(html_text, 'html.parser')
    trechos = []
    for nt in soup.select('div.noteText'):
        bruto = nt.get_text(" ", strip=True)
        limpo = _KINDLE_META.sub("", bruto).strip()
        if limpo:
            trechos.append("> " + limpo)
    return "\n\n".join(trechos)

def cli_kindle2md(html_path: str) -> None:
    src = Path(html_path)
    md_text = kindle_extract_html_to_md(src.read_text(encoding="utf-8"))
    destino = src.with_suffix(".md")
    destino.write_text(md_text, encoding="utf-8")
    print(f"Nota criada: {destino}")

# ============ Prompts (md2anki / txt2anki) ============
ANALYSIS_PROMPT = textwrap.dedent("""
You are a bilingual vocabulary analyst for learners of English whose native language is Portuguese.
You will receive a block of English text.  For each sentence or distinct clause in the text you must:

1. Provide the sentence in English exactly as it appears (use the key `sentence_en`).
2. Provide your best idiomatic translation of that sentence into Portuguese (use the key `sentence_pt`).
3. Identify every word or expression in that sentence that might be challenging for an intermediate learner of English.  These include uncommon nouns, verbs, adjectives, adverbs, phrasal verbs, collocations (such as compound nouns and adjective+noun phrases), and idiomatic expressions.  When a challenging item is part of a multiword expression or collocation (for example, a verb plus a particle or adverb like "held on tightly", or a compound noun like "knee‑high stockings"), return the full expression as the `word`.
    • `word`: the word or expression exactly as it appears in the sentence (lowercase).
    • `translation`: a concise Portuguese translation used in the sentence.
    • `definitions`: a list of objects. Each object must contain:
        – `definition`: an English definition.
        – `translations`: a list of possible Portuguese translations for that definition.
        – `contextual`: true if this definition matches the meaning in the sentence, otherwise false.

Return the result strictly as a JSON array.  Each element of the array must be an object with exactly three keys: `sentence_en`, `sentence_pt` and `terms` (an array of term objects described above).  Do not wrap the JSON in any code fence, do not add comments, and do not include any keys other than those requested.  The JSON must be valid and parsable by standard JSON parsers.

Text to analyse:
{md_text}
""")

PROMPT_CARDS = textwrap.dedent("""
You are an advanced English vocabulary tutor.

Produce 3 sentences (present, past, future) that NATURALLY illustrate the WORD
with the meaning suggested by its Portuguese translation.

OUTPUT FORMAT (strict):
S_PRESENT: <English sentence in PRESENT SIMPLE containing «WORD»>
PT_PRESENT: <idiomatic PT translation of S_PRESENT>
S_PAST: <English sentence in PAST SIMPLE containing «WORD»>
PT_PAST: <idiomatic PT translation of S_PAST>
S_FUTURE: <English sentence in FUTURE SIMPLE (will) containing «WORD»>
PT_FUTURE: <idiomatic PT translation of S_FUTURE>
---BLOCK---
1. Pronúncia
IPA: /.../ — Soa como: <dica PT>

2. Significado
Definição curta: [descrição PT de 1 linha]

3. Fato curioso
<curiosity in Portuguese>  (omit if none)

4. Forma base + classe gramatical
Base: <span style="color:#00aaff;font-weight:bold;">[lemma_or_expression]</span><br>
Classe: [classe em PT]

Rules:
• Labels S_* / PT_* must appear EXACTLY as shown, ALL CAPS followed by a colon.
  Missing any label → respond with FAIL.
• If WORD is a verb in inflected form, use its lemma in the sentences,
  conjugated to each tense; otherwise keep unchanged.
• Wrap each occurrence of the target item with « … » in English sentences.
• Each English sentence ≤ 14 words.
• Return exactly the fields above, no extra commentary.

WORD: «{word}»
Portuguese context meaning: «{pt}»
""")

FALLBACK_PROMPT = textwrap.dedent("""
You are an English tutor helping Portuguese speakers learn vocabulary.
Given an expression and its Portuguese meaning, generate three concise
example sentences in English: one in the present simple, one in the past
simple, and one in the future simple (using "will").  Each sentence must
naturally illustrate the expression and the meaning provided.  After each
English sentence, provide the idiomatic Portuguese translation.  Use the
following labels exactly as shown:
PRESENT_EN: <English sentence>
PRESENT_PT: <Portuguese translation>
PAST_EN: <English sentence>
PAST_PT: <Portuguese translation>
FUTURE_EN: <English sentence>
FUTURE_PT: <Portuguese translation>
Return exactly these six lines and nothing else.

Expression: «{expression}»
Meaning in Portuguese: «{meaning}»
""")

# ============ Helpers de formatação ============
def block_html(raw: str, meaning_pt: Optional[str] = None) -> str:
    raw = re.sub(r'^\d\.\s.*', lambda m: f"<b>{m.group(0)}</b>", raw.strip(), flags=re.MULTILINE)
    raw = re.sub(r'\n(<b>\d\.\s)', r'\n<br>\1', raw)
    if meaning_pt:
        def repl(m):
            texto = re.sub(
                re.escape(meaning_pt),
                lambda mm: ("<span style='background-color:#ff9800;color:#fff;"
                            "font-weight:bold;border-radius:4px;padding:0 4px;'>"
                            + mm.group(0) + "</span>"),
                m.group(1),
                flags=re.IGNORECASE,
            )
            return f"Definição curta: {texto}"
        raw = re.sub(r'Definição curta:\s*(.*)', repl, raw, flags=re.IGNORECASE)
    return highlight_html(raw).replace("\n", "<br>")

def format_definitions(word: str, definitions: List[Dict[str, Any]]) -> str:
    lines = [f"Significados: {word}"]
    for idx, d in enumerate(definitions, 1):
        def_text = d.get('definition', '')
        trans = ', '.join(d.get('translations', []))
        line = f"{idx}. {def_text}"
        if trans:
            line += f" | {trans}"
        if d.get('contextual'):
            style = ("background-color:#ff9800;color:#fff;font-weight:bold;"
                     "border-radius:4px;padding:0 4px;")
            line = f"<span style=\"{style}\">{line}</span>"
            line = f"<b>{line}</b>"
        lines.append(line)
    return "<br>".join(lines)

# ============ Geração de cartões (com OpenAI) ============
def build_cards(word: str, meaning_pt: str, extra_translations: Optional[List[str]] = None) -> List[Tuple[str,str]]:
    """Gera pares (frente_html, verso_html). Pode lançar RuntimeError em caso de formato inválido."""
    for i in range(MAX_TRIES):
        temp = 0.2 if i < MAX_TRIES-1 else 0.0
        resp = get_client().chat.completions.create(
            model=CHAT_MODEL,
            messages=[{"role":"user","content":PROMPT_CARDS.format(word=word, pt=meaning_pt)}],
            temperature=temp
        ).choices[0].message.content

        if resp.strip().startswith("FAIL") or '---BLOCK---' not in resp:
            continue
        head, block = resp.split('---BLOCK---', 1)
        triples = PAIR.findall(head)
        # Aceita qualquer número de pares >=1; se retornou 3, ótimo.
        if not triples:
            continue
        html = block_html(block, meaning_pt)
        cards = []
        for _, en, pt_sent in triples:
            pt_highlight = re.sub(
                re.escape(meaning_pt),
                lambda mm: ("<span style='color:#ff9800;font-weight:bold;'>"
                            + mm.group(0) + "</span>"),
                pt_sent.strip(),
                flags=re.IGNORECASE,
            )
            back = highlight_html(pt_highlight) + "<br><br>" + html
            if extra_translations:
                back += "<br><br><b>Possíveis traduções:</b> " + ", ".join(extra_translations)
            cards.append((highlight_html(en.strip()), back))
        return cards
    raise RuntimeError("formato inválido")

def fallback_sentences(expression: str, meaning: str) -> List[Tuple[str,str]]:
    prompt = FALLBACK_PROMPT.format(expression=expression, meaning=meaning)
    try:
        resp = get_client().chat.completions.create(
            model=CHAT_MODEL, messages=[{"role":"user","content":prompt}], temperature=0.2
        ).choices[0].message.content
    except Exception:
        return []
    lines = resp.strip().splitlines()
    pairs = []
    cur = None
    for ln in lines:
        ln = ln.strip()
        if ln.startswith("PRESENT_EN:"):
            cur = ln[len("PRESENT_EN:"):].strip()
        elif ln.startswith("PRESENT_PT:") and cur:
            pairs.append((cur, ln[len("PRESENT_PT:"):].strip())); cur=None
        elif ln.startswith("PAST_EN:"):
            cur = ln[len("PAST_EN:"):].strip()
        elif ln.startswith("PAST_PT:") and cur:
            pairs.append((cur, ln[len("PAST_PT:"):].strip())); cur=None
        elif ln.startswith("FUTURE_EN:"):
            cur = ln[len("FUTURE_EN:"):].strip()
        elif ln.startswith("FUTURE_PT:") and cur:
            pairs.append((cur, ln[len("FUTURE_PT:"):].strip())); cur=None
    return pairs

# ============ Análise de Markdown (md2anki) ============
def analyse_md_content(content: str) -> List[Dict[str, Any]]:
    snippet = content if len(content) <= 12000 else content[:12000]
    prompt = ANALYSIS_PROMPT.format(md_text=snippet)
    response = get_client().chat.completions.create(
        model=CHAT_MODEL, messages=[{"role":"user","content":prompt}], temperature=0.2
    ).choices[0].message.content
    try:
        data = json.loads(response)
    except Exception as e:
        raise RuntimeError("Failed to parse JSON from analysis: " + str(e))
    if not isinstance(data, list):
        raise RuntimeError("Analysis result is not a list")
    # normalização
    for item in data:
        terms = item.get('terms', [])
        if isinstance(terms, list):
            for t in terms:
                defs = t.get('definitions', [])
                if isinstance(defs, list):
                    normalized = []
                    for d in defs:
                        if isinstance(d, dict):
                            tr = d.get('translations', [])
                            if isinstance(tr, str):
                                translations = [x.strip() for x in tr.split(',') if x.strip()]
                            elif tr is None:
                                translations = []
                            elif isinstance(tr, list):
                                translations = [str(x).strip() for x in tr if str(x).strip()]
                            else:
                                translations = [str(tr).strip()]
                            normalized.append({
                                'definition': d.get('definition','').strip(),
                                'translations': translations,
                                'contextual': bool(d.get('contextual', False))
                            })
                        elif isinstance(d, str):
                            normalized.append({
                                'definition': d.strip(),
                                'translations': [],
                                'contextual': False
                            })
                    t['definitions'] = normalized
                else:
                    t['definitions'] = []
    return data

def create_pdf_from_analysis(items: List[Dict[str, Any]], pdf_path: str) -> None:
    """Gera um PDF via PPTX + LibreOffice; cai para PPTX se LibreOffice não existir."""
    if Presentation is None:
        raise RuntimeError("python-pptx não está instalado.")
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, "analysis_report.pptx")
        prs = Presentation()
        blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
        for idx, item in enumerate(items, 1):
            slide = prs.slides.add_slide(blank)
            shapes = slide.shapes
            # título
            tb = shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(1.2))
            tf = tb.text_frame; tf.clear()
            p = tf.paragraphs[0]; run = p.add_run(); run.font.size = Pt(18); run.font.bold = True
            sentence_en = item.get('sentence_en','').replace('\n',' ').strip()
            run.text = f'{idx}. "{sentence_en}"'
            # corpo
            body = shapes.add_textbox(Inches(0.3), Inches(1.6), Inches(9.4), Inches(3.8)).text_frame
            body.clear()
            para = body.paragraphs[0]
            para.text = "Tradução: " + item.get('sentence_pt','')
            para.font.size = Pt(14)
            terms = item.get('terms', [])
            for term in terms:
                q = body.add_paragraph()
                word = term.get('word','')
                defs = term.get('definitions', [])
                parts = []
                for d in defs:
                    line = d.get('definition','')
                    tr = ', '.join(d.get('translations', []))
                    if tr: line += f" | {tr}"
                    if d.get('contextual'): line = f"*{line}*"
                    parts.append(line)
                q.text = f"{word} | {'; '.join(parts) if parts else ''}"
                q.font.size = Pt(12)
        prs.save(pptx_path)
        try:
            subprocess.run(["libreoffice","--headless","--convert-to","pdf","--outdir",tmpdir,pptx_path],
                           check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            generated_pdf = os.path.join(tmpdir, "analysis_report.pdf")
            # Alguns LibreOffice nomeiam com base — garantimos o caminho
            if not os.path.exists(generated_pdf):
                base = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
                generated_pdf = os.path.join(tmpdir, base)
            with open(generated_pdf,'rb') as src, open(pdf_path,'wb') as dst:
                dst.write(src.read())
        except FileNotFoundError:
            fallback_pptx = pdf_path.replace(".pdf",".pptx")
            with open(pptx_path,'rb') as src, open(fallback_pptx,'wb') as dst:
                dst.write(src.read())
            raise RuntimeError("LibreOffice não encontrado. Relatório PPTX salvo em " + fallback_pptx)
        except subprocess.CalledProcessError as e:
            raise RuntimeError("Falha ao converter PPTX para PDF: " + e.stderr.decode(errors='ignore'))

# ============ Pipelines principais ============
def run_txt2anki(txt_path: str) -> None:
    deck_setup()
    total: List[Tuple[str,str]] = []
    bad: List[str] = []
    with open(txt_path, encoding="utf-8") as f:
        for ln in f:
            if "|" not in ln: 
                continue
            w, t_raw = [p.strip() for p in ln.split("|",1)]
            meaning, extras = t_raw, []
            if "(" in t_raw and t_raw.endswith(")"):
                base, rest = t_raw.rsplit("(",1)
                meaning = base.strip()
                extra_str = rest[:-1].strip()
                extra_str = re.sub(r"^(poss[ií]veis tradu[cç][ãa]oes?:?)\s*", "", extra_str, flags=re.IGNORECASE)
                extras = [x.strip() for x in extra_str.split(",") if x.strip()]
            try:
                total += build_cards(w, meaning, extras)
            except Exception as e:
                print("⚠", w, "-", e)
                bad.append(w)
    if total:
        notes=[]
        for fr, bk in total:
            mp3 = tts_mp3(re.sub(r"<[^>]+>","", fr))
            notes.append({
                "deckName": DECK_NAME,
                "modelName": MODEL_NAME,
                "fields": {FIELDS[0]: fr + f"<br>[sound:{mp3}]", FIELDS[1]: bk},
                "options": {"allowDuplicate": False},
                "tags": ["auto_gpt","tts"]
            })
        anki("addNotes", notes=notes)
        print(f"► {len(total)} cartões enviados")
    if bad:
        print("\nFalharam:", ", ".join(bad))
    else:
        print("Tudo certo!")

def run_md2anki(md_file: str, filter_file: Optional[str] = None) -> None:
    deck_setup()
    content = Path(md_file).read_text(encoding='utf-8')
    analysis = analyse_md_content(content)
    # filtro opcional
    filter_set = set()
    if filter_file and os.path.exists(filter_file):
        for line in Path(filter_file).read_text(encoding='utf-8').splitlines():
            t = line.strip()
            if t: filter_set.add(t.lower())
    elif filter_file:
        print(f"Arquivo de filtro '{filter_file}' não encontrado; ignorando.")
    # mapa de definições
    def_map: Dict[Tuple[str,str], Dict[str, Any]] = {}
    for item in analysis:
        for term in item.get('terms', []):
            w = term.get('word','').strip().lower()
            t = term.get('translation','').strip().lower()
            if w and t:
                def_map[(w,t)] = {'definitions': term.get('definitions', [])}

    # pares únicos preservando ordem
    pairs: List[Tuple[str,str]] = []
    seen = set()
    for item in analysis:
        for term in item.get('terms', []):
            word = term.get('word','').strip()
            trans = term.get('translation','').strip()
            key = (word.lower(), trans.lower())
            if word and trans and key not in seen:
                include = False
                if not filter_set:
                    include = True
                else:
                    wl = word.lower()
                    for flt in filter_set:
                        if flt in wl or wl in flt:
                            include = True; break
                if include:
                    pairs.append((word, trans)); seen.add(key)

    built_cards: List[Tuple[str,str]] = []
    fallback_pairs: List[Tuple[str,str]] = []
    for word, translation in pairs:
        try:
            built_cards += build_cards(word, translation)
        except Exception:
            print(f"⚠ {word} – cartão simples será gerado")
            fallback_pairs.append((word, translation))

    # Envia os cartões principais (um a um para lidar com duplicados)
    added = 0
    for front, back in built_cards:
        mp3 = tts_mp3(re.sub(r"<[^>]+>", "", front))
        note = {
            "deckName": DECK_NAME,
            "modelName": MODEL_NAME,
            "fields": {FIELDS[0]: front + f"<br>[sound:{mp3}]", FIELDS[1]: back},
            "options": {"allowDuplicate": False},
            "tags": ["auto_gpt","tts"]
        }
        try:
            anki("addNote", note=note); added += 1
        except RuntimeError as e:
            if 'duplicate' in str(e).lower(): 
                continue
            else:
                raise
    # Fallback
    for word, translation in fallback_pairs:
        key = (word.strip().lower(), translation.strip().lower())
        info = def_map.get(key, {})
        definitions = info.get('definitions', [])
        def_html = format_definitions(word, definitions) if definitions else ""
        pairs = fallback_sentences(word, translation) or []
        for en, pt in pairs:
            mp3 = tts_mp3(en)
            front_text = highlight_html(en) + f"<br>[sound:{mp3}]"
            back_text  = highlight_html(pt) + ("<br><br>"+def_html if def_html else "")
            note = {
                "deckName": DECK_NAME,
                "modelName": MODEL_NAME,
                "fields": {FIELDS[0]: front_text, FIELDS[1]: back_text},
                "options": {"allowDuplicate": False},
                "tags": ["auto_gpt","tts","fallback"]
            }
            try:
                anki("addNote", note=note); added += 1
            except RuntimeError as e:
                if 'duplicate' in str(e).lower(): 
                    continue
                else:
                    raise
    print(f"► {added} cartões enviados")

    # Relatório PDF (ou PPTX se LibreOffice indisponível)
    base = os.path.splitext(os.path.basename(md_file))[0]
    pdf_path = os.path.join(os.path.dirname(md_file), base + "_study.pdf")
    try:
        create_pdf_from_analysis(analysis, pdf_path)
        print("Relatório PDF salvo em:", pdf_path)
    except RuntimeError as e:
        print(str(e))

# ============ URL -> estudo -> Anki ============
def extract_text_from_url(url: str) -> str:
    if BeautifulSoup is None:
        raise RuntimeError("beautifulsoup4 não está instalado.")
    resp = requests.get(url)
    soup = BeautifulSoup(resp.content, 'html.parser')
    for tag in soup(['script','style']):
        tag.decompose()
    texto = soup.get_text(separator='\n')
    linhas = [l.strip() for l in texto.splitlines() if l.strip()]
    return '\n'.join(linhas)

def generate_verb_forms(expr: str) -> List[str]:
    parts = expr.split()
    nlp = get_nlp()
    if nlp is not None and parts:
        lemma = nlp(parts[0])[0].lemma_
    else:
        # fallback simples
        lemma = parts[0] if parts else expr
    forms = [lemma, lemma + 's', lemma + 'ed', lemma + 'ing']
    if len(parts) == 2:
        return [f + ' ' + parts[1] for f in forms]
    return forms

def structure_didactic_from_raw(content: str) -> List[Dict[str, Any]]:
    """Converte conteúdo cru de uma página em estrutura didática (lista de entradas)."""
    prompt = ("Por favor, responda APENAS com o JSON no formato exato solicitado, sem texto adicional."
              + textwrap.dedent("""
Você é um especialista em ensino de inglês para brasileiros.
Abaixo está o conteúdo bruto de um site explicando o uso de uma palavra, expressão ou phrasal verb.
Sua tarefa é:
1. Extrair o lema (palavra ou expressão estudada).
2. Identificar todas as definições distintas.
3. Para cada definição, retornar:
   - classe_gramatical (verbo, substantivo etc.)
   - pronúncia (IPA)
   - definicao (explicação em português)
   - curiosidades (se houver)
   - frases (lista com 3 objetos para tempos presente, passado, futuro):
       {"tempo":"presente","ingles":"...","portugues":"..."},
       {"tempo":"passado","ingles":"...","portugues":"..."},
       {"tempo":"futuro","ingles":"...","portugues":"..."}

Formato JSON estrito:
[
  {
    "expressao": "...",
    "classe_gramatical": "...",
    "pronuncia": "...",
    "definicao": "...",
    "curiosidades": "...",
    "frases": [
      {"tempo":"presente","ingles":"...","portugues":"..."},
      {"tempo":"passado","ingles":"...","portugues":"..."},
      {"tempo":"futuro","ingles":"...","portugues":"..."}
    ]
  },
  ...
]

Conteúdo:
{conteudo}
""").format(conteudo=content))
    last_err = None
    for i in range(MAX_TRIES):
        try:
            resp = get_client().chat.completions.create(
                model=CHAT_MODEL, messages=[{"role":"user","content":prompt}], temperature=0.3 if i<MAX_TRIES-1 else 0.0
            ).choices[0].message.content.strip()
            raw = re.sub(r'^```(?:json)?\n?', '', resp)
            raw = re.sub(r'\n?```$', '', raw)
            data = json.loads(raw)
            if isinstance(data, list):
                return data
        except Exception as e:
            last_err = e
            prompt = ("Formato incorreto, reimprima apenas o JSON solicitado."
                      + prompt.split("Conteúdo:")[0] + "Conteúdo:\n" + content)
    raise RuntimeError("GPT não retornou JSON válido após várias tentativas: " + str(last_err))

def send_structured_cards_to_anki(cards: List[Dict[str, Any]]) -> int:
    """Envia notas geradas por structure_didactic_from_raw ao Anki (com áudio)."""
    total = 0
    defs_html = "<br>".join(f"{i+1}. {c.get('definicao','')}" for i, c in enumerate(cards))
    for entry in cards:
        expr = entry['expressao']
        forms = generate_verb_forms(expr)
        pattern = rf"\\b(?:{'|'.join(re.escape(f) for f in forms)})\\b"
        defs_block = f"<span style='color:#00aaff;font-weight:bold;'>{expr}</span><br>{defs_html}"
        for frm in entry['frases']:
            audio_data = get_client().audio.speech.create(model=TTS_MODEL, voice=random.choice(VOICES), input=frm['ingles']).content
            audio_file = f"tts_{hashlib.sha1(audio_data).hexdigest()[:10]}.mp3"
            anki("storeMediaFile", filename=audio_file, data=base64.b64encode(audio_data).decode())

            destaque_en = re.sub(pattern,
                                 lambda m: f"<span style='color:#00aaff;font-weight:bold;'>{m.group(0)}</span>",
                                 frm['ingles'], flags=re.IGNORECASE)
            destaque_pt = re.sub(re.escape(entry['definicao']),
                                 lambda m: f"<span style='color:#00aaff;font-weight:bold;'>{m.group(0)}</span>",
                                 frm['portugues'], flags=re.IGNORECASE)
            fr_text = f"{destaque_en}<br>[sound:{audio_file}]"
            bk_text = (
                f"{destaque_pt}<br>"
                f"<b>1. Pronúncia</b><br>"
                f"IPA: {entry.get('pronuncia','-')}<br><br>"
                f"<b>2. Significado</b><br>"
                f"{defs_block}<br><br>"
                f"<b>3. Fato curioso</b><br>"
                f"{entry.get('curiosidades','-')}<br><br>"
                f"<b>4. Forma base + classe gramatical</b><br>"
                f"Base: <span style='color:#00aaff;font-weight:bold;'>{expr}</span><br>"
                f"Classe: {entry.get('classe_gramatical','-')}"
            )
            note = {
                "deckName": DECK_NAME,
                "modelName": MODEL_NAME,
                "fields": {FIELDS[0]: fr_text, FIELDS[1]: bk_text},
                "options": {"allowDuplicate": False},
                "tags": ["auto_gpt","tts"]
            }
            try:
                anki("addNote", note=note)
                total += 1
            except RuntimeError as e:
                if 'duplicate' in str(e).lower():
                    continue
                else:
                    raise
    return total

def run_url2anki(url: str) -> None:
    texto = extract_text_from_url(url)
    cards = structure_didactic_from_raw(texto)
    n = send_structured_cards_to_anki(cards)
    print(f"► {n} cartões enviados.")

# ============ Embutir index.html e escrever se desejado ============
INDEX_HTML = r"""<!DOCTYPE html>
<html lang="pt">
<head>
  <meta charset="UTF-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>Learn English - Upload</title>
  <!-- Conteúdo original do seu index.html foi sintetizado aqui; 
       se preferir, substitua por sua versão integral. -->
</head>
<body>
  <h1>Learn English - Upload</h1>
  <p>Esta página é um placeholder. Use o CLI para processar arquivos (kindle2md, txt2anki, md2anki, url2anki).</p>
</body>
</html>
"""

def cli_write_ui(out_path: str) -> None:
    """Escreve um index.html simples. Pode substituir o conteúdo pela sua versão integral."""
    Path(out_path).write_text(INDEX_HTML, encoding="utf-8")
    print("index.html escrito em:", out_path)

# ============ CLI dispatcher ============
def main(argv: List[str]) -> None:
    if len(argv) < 2:
        print("""Uso:
  python anki_suite.py kindle2md arquivo.html
  python anki_suite.py txt2anki palavras.txt
  python anki_suite.py md2anki arquivo.md [filtro.txt]
  python anki_suite.py url2anki https://exemplo.com/estudo
  python anki_suite.py write-ui caminho/saida/index.html
""")
        sys.exit(1)
    cmd = argv[1].lower()
    try:
        if cmd == "kindle2md":
            cli_kindle2md(argv[2])
        elif cmd == "txt2anki":
            run_txt2anki(argv[2])
        elif cmd == "md2anki":
            md = argv[2]
            filt = argv[3] if len(argv) >= 4 else None
            run_md2anki(md, filt)
        elif cmd == "url2anki":
            run_url2anki(argv[2])
        elif cmd == "write-ui":
            cli_write_ui(argv[2])
        else:
            print("Comando não reconhecido:", cmd)
            sys.exit(2)
    except IndexError:
        print("Parâmetros insuficientes para o comando:", cmd)
        sys.exit(2)

if __name__ == "__main__":
    main(sys.argv)
